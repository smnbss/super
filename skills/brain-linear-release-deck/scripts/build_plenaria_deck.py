#!/usr/bin/env python3
"""Build a WeRoad "Shipping the Bets" Plenaria Tech deck from a JSON spec.

Visual template reverse-engineered from
`2026 05 - Plenaria - May 2026 - Tech` (Drive ID
1KIfGGs11BBS-Q8QHpSTh_bH9b3vuX8RWhUDvgahoteg).

The model writes a deck_spec.json (see references/deck-spec-schema.md) and
this script handles every coordinate, font, and color decision.
"""

from __future__ import annotations

import argparse
import json
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------- Brand ----------

WR_RED = RGBColor(0xFF, 0x47, 0x58)
WR_BLACK = RGBColor(0x00, 0x00, 0x00)
WR_DARK_GREY = RGBColor(0x4D, 0x4D, 0x4D)
WR_MID_GREY = RGBColor(0x9E, 0x9E, 0x9E)
WR_LIGHT_GREY = RGBColor(0xE5, 0xE5, 0xE5)
WR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_TITLE = "Poppins"           # ExtraBold weight comes from bold=True
FONT_BODY = "Poppins"
FONT_MONO = "Roboto"

# Slide is 10 × 5.625 in (16:9 widescreen)
SLIDE_W_IN = 10.0
SLIDE_H_IN = 5.625

# Logo placement — bottom-left, ~0.35in tall
LOGO_LEFT = Inches(0.30)
LOGO_TOP = Inches(5.10)
LOGO_W = Inches(1.10)
LOGO_H = Inches(0.35)


# ---------- Resource resolution ----------

def find_resources_dir() -> Path | None:
    """Locate the weroad-presentations skill's resources directory.

    We don't bundle logos in this skill — we reuse the canonical brand assets
    from weroad-presentations. Search likely locations.
    """
    candidates = [
        Path.home() / ".claude/plugins/weroad-ai/plugins/presentations/skills/weroad-presentations/resources",
        Path.home() / ".claude/plugins/cache/weroad-ai/presentations/1.0.0/skills/weroad-presentations/resources",
        Path.home() / ".claude/plugins/marketplaces/weroad-ai/plugins/presentations/skills/weroad-presentations/resources",
    ]
    for c in candidates:
        if c.is_dir() and (c / "WeRoad Logos").is_dir():
            return c
    return None


# ---------- Drawing helpers ----------

def add_textbox(slide, left, top, width, height, text, *,
                font=FONT_BODY, size=14, color=WR_BLACK,
                bold=False, italic=False, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, line_spacing=1.15,
                word_wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
    return box


def add_run(paragraph, text, *, font=FONT_BODY, size=14,
            color=WR_BLACK, bold=False, italic=False):
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return run


def add_eyebrow(slide, label, *, top=Inches(0.55), color=WR_RED):
    """Small red horizontal line + uppercase eyebrow label below it."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), top, Inches(0.45), Emu(38100)
    )
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = color
    add_textbox(
        slide, Inches(0.5), top + Inches(0.12), Inches(6), Inches(0.3),
        label.upper(), font=FONT_TITLE, size=11, color=color, bold=True,
    )


def add_logo(slide, dark_bg=False):
    res = find_resources_dir()
    if res is None:
        return  # logo missing → skip silently rather than fail
    logos = res / "WeRoad Logos"
    fname = "WeRoad_logo-w.png" if dark_bg else "WeRoad_logo.png"
    path = logos / fname
    if not path.is_file():
        return
    slide.shapes.add_picture(str(path), LOGO_LEFT, LOGO_TOP, height=LOGO_H)


def fill_slide_background(slide, color):
    """Paint the entire slide a flat color by adding a full-bleed rectangle."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W_IN), Inches(SLIDE_H_IN)
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    # Send to back so other shapes draw on top
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


# ---------- Slide builders ----------

def build_cover(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill_slide_background(slide, WR_BLACK)

    # Eyebrow line + label (white on black)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.30), Inches(0.45), Emu(38100)
    )
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = WR_RED
    add_textbox(
        slide, Inches(0.5), Inches(1.42), Inches(6), Inches(0.3),
        spec.get("cover_eyebrow", "PRODUCT REVIEW").upper(),
        font=FONT_TITLE, size=12, color=WR_RED, bold=True,
    )

    # Big title — written per window in the spec. If the spec forgets one,
    # we fall back to a window-derived placeholder rather than a fixed
    # phrase, so missing titles are visually obvious instead of silently
    # reusing last month's cover line.
    cover_title = spec.get("cover_title") or (
        f"{spec.get('window_label', 'Tech')} Review"
    )
    add_textbox(
        slide, Inches(0.5), Inches(1.85), Inches(9), Inches(1.85),
        cover_title,
        font=FONT_TITLE, size=60, color=WR_WHITE, bold=True,
        line_spacing=1.0,
    )
    # Subtitle
    add_textbox(
        slide, Inches(0.5), Inches(3.85), Inches(9), Inches(0.5),
        spec.get("cover_subtitle", "AI · Conversion · Operations · Community"),
        font=FONT_TITLE, size=22, color=WR_WHITE,
    )
    # Source line
    add_textbox(
        slide, Inches(0.5), Inches(4.45), Inches(9), Inches(0.3),
        spec.get("cover_source", "Source: Linear release notes"),
        font=FONT_TITLE, size=11, color=WR_MID_GREY, italic=True,
    )

    add_logo(slide, dark_bg=True)


def build_by_numbers(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_eyebrow(slide, "BY THE NUMBERS")

    # Auto-shrink the headline + push cards down when copy wraps — long
    # window-specific headlines are valuable, so the layout adapts.
    headline = spec.get("headline_count", "Releases shipped.")
    if len(headline) > 50:
        h_size, card_top = 28, Inches(2.45)
    elif len(headline) > 32:
        h_size, card_top = 32, Inches(2.25)
    else:
        h_size, card_top = 36, Inches(2.10)
    add_textbox(
        slide, Inches(0.5), Inches(0.95), Inches(9), Inches(1.40),
        headline,
        font=FONT_TITLE, size=h_size, color=WR_BLACK, bold=True,
        line_spacing=1.10,
    )

    cards = spec.get("by_numbers", [])[:4]
    if not cards:
        add_logo(slide)
        return

    # 4 evenly-spaced rounded-rectangle cards in a row. Card height shrinks
    # when the headline pushed the row down, keeping the logo footer clear.
    n = len(cards)
    margin = Inches(0.5)
    gap = Inches(0.20)
    card_w = (Inches(SLIDE_W_IN) - margin * 2 - gap * (n - 1)) / n
    card_h = Inches(5.05) - card_top

    for i, card in enumerate(cards):
        left = margin + (card_w + gap) * i
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, card_top, card_w, card_h
        )
        shape.adjustments[0] = 0.10  # slightly rounded
        shape.fill.solid()
        shape.fill.fore_color.rgb = WR_WHITE
        shape.line.color.rgb = WR_LIGHT_GREY
        shape.line.width = Pt(1)

        # Big red number — auto-shrink for long strings so it stays in the card
        # Card width is ~2.1in at 4 cards; ~2.4in at 3. Tune sizes to that.
        number_text = card.get("number", "")
        n_len = len(number_text)
        if n_len >= 7:
            num_size = 32
        elif n_len >= 6:
            num_size = 40
        elif n_len >= 4:
            num_size = 52
        else:
            num_size = 64
        add_textbox(
            slide, left, card_top + Inches(0.25), card_w, Inches(1.0),
            number_text,
            font=FONT_TITLE, size=num_size, color=WR_RED, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        # Bold label
        add_textbox(
            slide, left + Inches(0.10), card_top + Inches(1.40),
            card_w - Inches(0.20), Inches(0.5),
            card.get("label", ""),
            font=FONT_TITLE, size=14, color=WR_BLACK, bold=True,
            align=PP_ALIGN.CENTER,
        )
        # Caption
        add_textbox(
            slide, left + Inches(0.10), card_top + Inches(1.85),
            card_w - Inches(0.20), Inches(0.5),
            card.get("caption", ""),
            font=FONT_TITLE, size=11, color=WR_MID_GREY,
            align=PP_ALIGN.CENTER,
        )

    add_logo(slide)


def build_four_tracks(prs, spec):
    """Slide 3 — overview columns, one per track. Accepts both `tracks`
    (preferred) and the legacy `bets` key for transitional specs."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_eyebrow(slide, "THE STORY")
    title_text = (spec.get("four_tracks_title")
                  or spec.get("four_bets_title")
                  or "Four tracks, one direction")
    # Auto-shrink + adjust column-top so a 2-line title doesn't bleed
    # into the columns below.
    if len(title_text) > 38:
        t_size, col_top = 28, Inches(1.95)
    elif len(title_text) > 28:
        t_size, col_top = 32, Inches(1.85)
    else:
        t_size, col_top = 36, Inches(1.80)
    add_textbox(
        slide, Inches(0.5), Inches(0.95), Inches(9), Inches(0.95),
        title_text,
        font=FONT_TITLE, size=t_size, color=WR_BLACK, bold=True,
        line_spacing=1.05,
    )

    tracks = (spec.get("tracks") or spec.get("bets") or [])[:4]
    n = len(tracks) or 4
    margin = Inches(0.5)
    gap = Inches(0.18)
    col_w = (Inches(SLIDE_W_IN) - margin * 2 - gap * (n - 1)) / n
    # Columns stop just above the logo footer (logo top is at 5.10in).
    col_h = Inches(5.05) - col_top

    for i, track in enumerate(tracks):
        left = margin + (col_w + gap) * i
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, col_top, col_w, col_h
        )
        shape.adjustments[0] = 0.06
        shape.fill.solid()
        shape.fill.fore_color.rgb = WR_WHITE
        shape.line.color.rgb = WR_LIGHT_GREY
        shape.line.width = Pt(1)

        # Header (red)
        add_textbox(
            slide, left + Inches(0.20), col_top + Inches(0.18),
            col_w - Inches(0.40), Inches(0.30),
            track.get("name", "").upper(),
            font=FONT_TITLE, size=12, color=WR_RED, bold=True,
        )
        # Claim — auto-shrink so window-specific multi-clause sentences
        # fit the ~2-inch column without overrunning into the bullet list.
        claim_text = track.get("claim", "")
        c_len = len(claim_text)
        if c_len > 70:
            claim_size = 10
        elif c_len > 50:
            claim_size = 11
        else:
            claim_size = 13
        # Claim slot leaves the bottom 2/3 of the column for bullets — most
        # of the column's vertical real estate goes to the release list.
        claim_top = col_top + Inches(0.55)
        claim_h = Inches(1.15)
        add_textbox(
            slide, left + Inches(0.20), claim_top,
            col_w - Inches(0.40), claim_h,
            claim_text,
            font=FONT_TITLE, size=claim_size, color=WR_BLACK, bold=True,
            line_spacing=1.10,
        )
        # Bullet list — capped at 4 (5+ overflows even with multi-line
        # wrap headroom). The bullet area takes everything from below the
        # claim down to the column bottom minus a small bottom margin.
        bullets = track.get("overview_bullets", [])[:4]
        bullet_top = claim_top + claim_h + Inches(0.10)
        bullet_h = (col_top + col_h) - bullet_top - Inches(0.15)
        bullet_text = "\n".join(f"·  {b}" for b in bullets)
        add_textbox(
            slide, left + Inches(0.20), bullet_top,
            col_w - Inches(0.40), bullet_h,
            bullet_text,
            font=FONT_TITLE, size=9, color=WR_BLACK, line_spacing=1.20,
        )

    add_logo(slide)


def build_chapter(prs, track, *, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_background(slide, WR_RED)

    # White eyebrow line + "TRACK N OF M · <name>"
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.30), Inches(0.45), Emu(38100)
    )
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = WR_WHITE
    name_part = (track.get("name") or "").upper()
    eyebrow_text = (f"TRACK {idx} OF {total}" if not name_part
                    else f"TRACK {idx} OF {total} · {name_part}")
    add_textbox(
        slide, Inches(0.5), Inches(1.42), Inches(8), Inches(0.3),
        eyebrow_text,
        font=FONT_TITLE, size=12, color=WR_WHITE, bold=True,
    )
    # Big title — auto-shrink based on length so 2-line titles don't collide
    chapter_title = track.get("chapter_title", track.get("name", ""))
    if len(chapter_title) >= 24:
        title_size, title_h = 48, Inches(2.0)
    elif len(chapter_title) >= 16:
        title_size, title_h = 56, Inches(1.85)
    else:
        title_size, title_h = 64, Inches(1.40)
    add_textbox(
        slide, Inches(0.5), Inches(1.85), Inches(9), title_h,
        chapter_title,
        font=FONT_TITLE, size=title_size, color=WR_WHITE, bold=True,
        line_spacing=1.0,
    )
    # Subtitle — positioned below the (variable) title height with 0.30in gap
    add_textbox(
        slide, Inches(0.5), Inches(1.85) + title_h + Inches(0.30),
        Inches(9), Inches(0.7),
        track.get("subtitle", ""),
        font=FONT_TITLE, size=18, color=WR_WHITE,
    )
    add_logo(slide, dark_bg=True)


def resolve_visual_path(content, assets_dir: Path | None) -> Path | None:
    """Resolve the image path for any slide type (content or spotlight).

    Order of precedence:
      1. explicit `visual.path` if set
      2. `<assets_dir>/<mol_id>.*` (content slides use this field)
      3. `<assets_dir>/<release_id>.*` (impact spotlights use this field)
      4. None → placeholder / no image
    Within an asset directory, GIF wins over still images when both exist.
    """
    visual = content.get("visual") or {}
    explicit = visual.get("path")
    if explicit:
        p = Path(explicit)
        return p if p.is_file() else None

    if assets_dir is None:
        return None

    issue_id = content.get("mol_id") or content.get("release_id")
    if not issue_id:
        return None

    candidates = sorted(assets_dir.glob(f"{issue_id}.*"))
    if not candidates:
        return None
    for p in candidates:
        if p.suffix.lower() == ".gif":
            return p
    return candidates[0]


def build_content(prs, content, assets_dir: Path | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_eyebrow(slide, content.get("eyebrow", ""))

    add_textbox(
        slide, Inches(0.5), Inches(0.95), Inches(9), Inches(1.2),
        content.get("title", ""),
        font=FONT_TITLE, size=32, color=WR_BLACK, bold=True,
        line_spacing=1.10,
    )

    # Body left, visual right
    body_left = Inches(0.5)
    body_top = Inches(2.20)
    body_w = Inches(4.7)
    body_h = Inches(2.80)

    bullets = content.get("bullets", [])
    body_box = slide.shapes.add_textbox(body_left, body_top, body_w, body_h)
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.30
        if i > 0:
            p.space_before = Pt(8)
        if isinstance(bullet, dict):
            lead = bullet.get("lead", "")
            body = bullet.get("body", "")
            if lead:
                add_run(p, f"{lead}  ", font=FONT_TITLE, size=12,
                        color=WR_RED, bold=True)
            if body:
                add_run(p, body, font=FONT_TITLE, size=12,
                        color=WR_BLACK)
        else:
            add_run(p, str(bullet), font=FONT_TITLE, size=12,
                    color=WR_BLACK)

    # Visual — resolve from explicit path, else from `mol_id` + assets dir
    visual = content.get("visual") or {}
    vis_left = Inches(5.40)
    vis_top = Inches(2.10)
    vis_w = Inches(4.10)
    vis_h = Inches(2.95)

    img_path = resolve_visual_path(content, assets_dir)
    if img_path is not None:
        # Let pptx auto-scale by width; centre vertically inside the slot
        # by computing the picture's natural ratio and adjusting top.
        pic = slide.shapes.add_picture(str(img_path), vis_left, vis_top,
                                       width=vis_w)
        if pic.height > vis_h:
            # Too tall when fit-to-width — re-add fitting height instead
            slide.shapes._spTree.remove(pic._element)
            slide.shapes.add_picture(str(img_path), vis_left, vis_top,
                                     height=vis_h)
    else:
        label = visual.get("label") or content.get("mol_id") or "visual placeholder"
        _placeholder(slide, vis_left, vis_top, vis_w, vis_h, label)

    add_logo(slide)


def _placeholder(slide, left, top, w, h, label):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    box.adjustments[0] = 0.04
    box.fill.solid()
    box.fill.fore_color.rgb = WR_LIGHT_GREY
    box.line.color.rgb = WR_MID_GREY
    box.line.width = Pt(0.75)
    add_textbox(slide, left, top, w, h, label,
                font=FONT_TITLE, size=11, color=WR_MID_GREY, italic=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def build_impact_spotlight(prs, spotlight, assets_dir: Path | None = None):
    """Hero slide for a single significant measured win — gives the deck
    a visual punch where the impact story is too important for a regular
    bullet list. Layout: huge red KPI number on the left, headline +
    context narrative on the right, attribution at the bottom.

    All fields are optional; the renderer adapts to what's present so
    minor wins (number + headline only) and big wins (full narrative
    with caveat + attribution) both look intentional.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_eyebrow(slide, "MEASURED IMPACT")

    kpi = spotlight.get("kpi_number", "")
    kpi_label = spotlight.get("kpi_label", "")
    kpi_period = spotlight.get("kpi_period", "")

    # Resolve the optional inset visual. When present, the left column
    # compresses the KPI to the upper third and gives the lower two
    # thirds to the image — readers see both the number and the product
    # change in one frame. When absent, the KPI uses the full left half
    # at the bigger sizes (the original text-only layout).
    img_path = resolve_visual_path(spotlight, assets_dir)

    kpi_left = Inches(0.4)
    kpi_box_w = Inches(4.7) if img_path else Inches(5.0)
    kpi_len = len(kpi)

    # Two size tables — smaller when sharing the column with an image.
    if img_path:
        if kpi_len >= 9:
            kpi_size = 44
        elif kpi_len >= 8:
            kpi_size = 50
        elif kpi_len >= 7:
            kpi_size = 60
        elif kpi_len >= 6:
            kpi_size = 68
        elif kpi_len >= 5:
            kpi_size = 80
        elif kpi_len >= 4:
            kpi_size = 96
        else:
            kpi_size = 116
        kpi_top = Inches(1.20)
        kpi_h = Inches(1.10)
        label_top = Inches(2.30)
    else:
        if kpi_len >= 9:
            kpi_size = 56
        elif kpi_len >= 8:
            kpi_size = 64
        elif kpi_len >= 7:
            kpi_size = 76
        elif kpi_len >= 6:
            kpi_size = 84
        elif kpi_len >= 5:
            kpi_size = 100
        elif kpi_len >= 4:
            kpi_size = 124
        else:
            kpi_size = 150
        kpi_top = Inches(1.55)
        kpi_h = Inches(2.20)
        label_top = Inches(3.85)

    add_textbox(
        slide, kpi_left, kpi_top, kpi_box_w, kpi_h,
        kpi,
        font=FONT_TITLE, size=kpi_size, color=WR_RED, bold=True,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0,
        word_wrap=False,
    )
    if kpi_label:
        add_textbox(
            slide, kpi_left, label_top, kpi_box_w, Inches(0.35),
            kpi_label,
            font=FONT_TITLE, size=13 if img_path else 14,
            color=WR_DARK_GREY, bold=True,
        )
    if kpi_period:
        add_textbox(
            slide, kpi_left, label_top + Inches(0.35),
            kpi_box_w, Inches(0.35),
            kpi_period,
            font=FONT_TITLE, size=10 if img_path else 11,
            color=WR_MID_GREY, italic=True,
        )

    # Inset visual — bottom-left quadrant. Sized to leave breathing
    # room above the logo footer (logo top at 5.10in).
    if img_path:
        inset_top = Inches(3.10)
        inset_left = Inches(0.4)
        inset_w = Inches(4.5)
        inset_h = Inches(1.85)
        try:
            pic = slide.shapes.add_picture(
                str(img_path), inset_left, inset_top, width=inset_w,
            )
            if pic.height > inset_h:
                slide.shapes._spTree.remove(pic._element)
                slide.shapes.add_picture(
                    str(img_path), inset_left, inset_top, height=inset_h,
                )
        except Exception:
            pass  # fall through silently if pptx can't embed the image

    # Right half: headline + context + caveat + attribution
    rx = Inches(5.55)
    rw = Inches(4.10)

    headline = spotlight.get("headline", "")
    h_len = len(headline)
    if h_len > 90:
        h_size, h_height = 18, Inches(1.55)
    elif h_len > 50:
        h_size, h_height = 22, Inches(1.30)
    else:
        h_size, h_height = 26, Inches(1.00)
    add_textbox(
        slide, rx, Inches(1.45), rw, h_height,
        headline,
        font=FONT_TITLE, size=h_size, color=WR_BLACK, bold=True,
        line_spacing=1.10,
    )

    # Context paragraph — keep its bottom edge above the caveat zone
    # (caveat starts at 4.30in). Context height is bounded so a long
    # paragraph clips/wraps inside its slot rather than bleeding over
    # the caveat below it.
    body_top = Inches(1.45) + h_height + Inches(0.15)
    caveat_top = Inches(4.30)
    ctx_max_h = caveat_top - body_top - Inches(0.10)
    context = spotlight.get("context", "")
    if context:
        add_textbox(
            slide, rx, body_top, rw, ctx_max_h,
            context,
            font=FONT_TITLE, size=11, color=WR_BLACK, line_spacing=1.25,
        )

    caveat = spotlight.get("caveat", "")
    if caveat:
        add_textbox(
            slide, rx, caveat_top, rw, Inches(0.65),
            caveat,
            font=FONT_TITLE, size=9, color=WR_MID_GREY, italic=True,
            line_spacing=1.20,
        )

    # Attribution along the bottom-right (release id + author when present)
    attr_parts = []
    rid = spotlight.get("release_id")
    rtitle = spotlight.get("release_title")
    author = spotlight.get("author")
    if rtitle:
        attr_parts.append(rtitle)
    if rid and rid not in (rtitle or ""):
        attr_parts.append(rid)
    if author:
        attr_parts.append(author)
    if attr_parts:
        add_textbox(
            slide, Inches(0.5), Inches(5.05), Inches(9.0), Inches(0.3),
            "  ·  ".join(attr_parts),
            font=FONT_TITLE, size=9, color=WR_MID_GREY,
            align=PP_ALIGN.RIGHT,
        )

    add_logo(slide)


def build_looking_ahead(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_eyebrow(slide, "LOOKING AHEAD")
    add_textbox(
        slide, Inches(0.5), Inches(0.95), Inches(9), Inches(0.7),
        spec.get("looking_ahead_title", "Where each track sharpens next"),
        font=FONT_TITLE, size=36, color=WR_BLACK, bold=True,
    )

    items = spec.get("looking_ahead", [])[:4]
    # 2x2 grid
    margin_left = Inches(0.5)
    grid_top = Inches(2.20)
    cell_w = Inches(4.4)
    cell_h = Inches(1.40)
    gap_x = Inches(0.20)
    gap_y = Inches(0.20)

    for i, item in enumerate(items):
        row = i // 2
        col = i % 2
        left = margin_left + (cell_w + gap_x) * col
        top = grid_top + (cell_h + gap_y) * row
        # Theme header (red)
        add_textbox(
            slide, left, top, cell_w, Inches(0.35),
            item.get("theme", "").upper(),
            font=FONT_TITLE, size=12, color=WR_RED, bold=True,
        )
        # Body
        add_textbox(
            slide, left, top + Inches(0.40), cell_w, cell_h - Inches(0.40),
            item.get("body", ""),
            font=FONT_TITLE, size=13, color=WR_BLACK, line_spacing=1.30,
        )

    add_logo(slide)


def build_thank_you(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_background(slide, WR_BLACK)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.95), Inches(0.45), Emu(38100)
    )
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = WR_RED

    add_textbox(
        slide, Inches(0.5), Inches(2.20), Inches(9), Inches(1.4),
        spec.get("thank_you_title", "Thank you."),
        font=FONT_TITLE, size=72, color=WR_WHITE, bold=True,
    )
    add_textbox(
        slide, Inches(0.5), Inches(3.65), Inches(9), Inches(0.6),
        spec.get("thank_you_subtitle",
                 "Questions, pushback, requests for deep dives — welcome."),
        font=FONT_TITLE, size=18, color=WR_WHITE,
    )
    add_logo(slide, dark_bg=True)


# ---------- Main ----------

def build(spec_path: Path, output_path: Path,
          assets_dir: Path | None = None) -> Path:
    spec = json.loads(spec_path.read_text())

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    # 1. Cover
    build_cover(prs, spec)
    # 2. By the numbers
    build_by_numbers(prs, spec)
    # 3. Four tracks, one direction
    build_four_tracks(prs, spec)
    # 3.5 (optional). Top-level impact spotlights — placed before the
    # first chapter so the audience sees the headline measured wins
    # before the bet-by-bet walk-through. Per-track spotlights live
    # inside that track's `slides` list and render in place.
    for spotlight in spec.get("impact_spotlights", []):
        build_impact_spotlight(prs, spotlight, assets_dir=assets_dir)
    # 4..N. Chapter + content slides per track
    tracks = spec.get("tracks") or spec.get("bets") or []
    total = len(tracks)
    for idx, track in enumerate(tracks, start=1):
        build_chapter(prs, track, idx=idx, total=total)
        for slide_spec in track.get("slides", []):
            stype = slide_spec.get("type", "content")
            if stype == "impact_spotlight":
                build_impact_spotlight(prs, slide_spec, assets_dir=assets_dir)
            else:
                build_content(prs, slide_spec, assets_dir=assets_dir)
    # Last. Thank you (no "looking ahead" slide — Plenaria reports what
    # shipped, doesn't preview what's coming and never makes asks of the
    # audience)
    build_thank_you(prs, spec)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--spec", required=True, type=Path,
                    help="path to deck_spec.json")
    ap.add_argument("--output", required=True, type=Path,
                    help="path to write the .pptx")
    ap.add_argument("--assets-dir", type=Path, default=None,
                    help="directory containing <mol_id>.{gif,png,jpg} images "
                         "fetched by fetch_linear_images.py")
    args = ap.parse_args()

    if not args.spec.is_file():
        print(f"spec not found: {args.spec}", file=sys.stderr)
        sys.exit(1)

    out = build(args.spec, args.output, assets_dir=args.assets_dir)
    print(f"wrote {out} ({out.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main()
